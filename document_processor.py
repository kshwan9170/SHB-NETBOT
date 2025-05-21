import os
from pathlib import Path
import tempfile
import re
import uuid
from typing import List, Dict, Any, Optional, Tuple

# Document processing libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    
try:
    import pandas as pd
except ImportError:
    pd = None

def process_document(file_path: str) -> List[Dict[str, Any]]:
    """
    Process a document file and extract text chunks with metadata
    
    Args:
        file_path: Path to the document file
        
    Returns:
        List of dictionaries containing text chunks and metadata
    """
    file_extension = Path(file_path).suffix.lower()
    filename = os.path.basename(file_path)
    doc_id = str(uuid.uuid4())
    
    chunks = []
    
    # Check for procedure guide Excel file and extract additional metadata
    is_procedure_guide = False
    guide_version = 'latest'
    
    # For Excel files, check if it's a procedure guide and extract version info
    if (file_extension == '.xlsx' or file_extension == '.xls') and \
       any(keyword in filename for keyword in ['업무 안내', '업무_안내', '업무안내', '가이드', '매뉴얼', '절차']):
        is_procedure_guide = True
        
        # Extract date pattern from filename (e.g., 업무 안내 가이드_2025.05.19.xlsx)
        date_pattern = re.search(r'_(\d{4}[.년\-_]\d{1,2}[.월\-_]\d{1,2})', filename)
        if date_pattern:
            guide_version = date_pattern.group(1)
    
    # Extract text based on file type
    if file_extension == '.pdf':
        raw_chunks = extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        raw_chunks = extract_text_from_docx(file_path)
    elif file_extension == '.pptx':
        raw_chunks = extract_text_from_pptx(file_path)
    elif file_extension == '.xlsx' or file_extension == '.xls':
        raw_chunks = extract_text_from_excel(file_path)
    elif file_extension == '.csv':
        raw_chunks = extract_text_from_csv(file_path)
    elif file_extension == '.txt':
        raw_chunks = extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {file_extension}")
    
    # Add metadata to each chunk
    for i, chunk in enumerate(raw_chunks):
        # Setup basic metadata
        metadata = {
            "filename": filename,
            "file_type": file_extension[1:],  # Remove the dot
            "chunk_index": i
        }
        
        # Enhanced metadata for procedure guide Excel files
        if is_procedure_guide:
            # Try to extract information about the Excel row if it's contained in the chunk
            if chunk.startswith("[업무 안내]"):
                # Parse sheet name from the source
                source_match = re.search(r'출처: .+ - (.+) 시트', chunk)
                if source_match:
                    metadata["sheet_name"] = source_match.group(1)
                
                # Extract row fields for better searchability
                for field in ["업무 유형", "질문 예시", "요약 응답", "상세 안내", "키워드"]:
                    field_match = re.search(f"{field}: ([^|]+)", chunk)
                    if field_match:
                        field_value = field_match.group(1).strip()
                        metadata[field] = field_value
                
                # Mark as procedure guide for special handling in retrieval
                metadata["content_type"] = "procedure_guide"
                metadata["guide_version"] = guide_version
        
        chunks.append({
            "doc_id": doc_id,
            "chunk_id": f"{doc_id}-{i}",
            "text": chunk,
            "metadata": metadata
        })
    
    return chunks

def extract_text_from_pdf(file_path: str) -> List[str]:
    """Extract text from PDF files"""
    if PyPDF2 is None:
        raise ImportError("PyPDF2 is required for PDF processing")
    
    text_chunks = []
    
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                
                if page_text:
                    # Clean and chunk the text
                    chunks = chunk_text(page_text)
                    for chunk in chunks:
                        if len(chunk.strip()) > 20:  # Only keep meaningful chunks
                            text_chunks.append(chunk)
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    
    return text_chunks

def extract_text_from_docx(file_path: str) -> List[str]:
    """Extract text from DOCX files"""
    if docx is None:
        raise ImportError("python-docx is required for DOCX processing")
    
    text_chunks = []
    
    try:
        doc = docx.Document(file_path)
        full_text = ""
        
        for para in doc.paragraphs:
            if para.text.strip():
                full_text += para.text + "\n"
        
        # Process tables separately to preserve structure
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    full_text += row_text + "\n"
        
        chunks = chunk_text(full_text)
        for chunk in chunks:
            if len(chunk.strip()) > 20:
                text_chunks.append(chunk)
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
    
    return text_chunks

def extract_text_from_pptx(file_path: str) -> List[str]:
    """Extract text from PPTX files"""
    if Presentation is None:
        raise ImportError("python-pptx is required for PPTX processing")
    
    text_chunks = []
    
    try:
        prs = Presentation(file_path)
        full_text = ""
        
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                full_text += "--- Slide ---\n" + slide_text + "\n"
        
        chunks = chunk_text(full_text)
        for chunk in chunks:
            if len(chunk.strip()) > 20:
                text_chunks.append(chunk)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    
    return text_chunks

def extract_text_from_txt(file_path: str) -> List[str]:
    """Extract text from TXT files"""
    text_chunks = []
    
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
            
            if content:
                chunks = chunk_text(content)
                for chunk in chunks:
                    if len(chunk.strip()) > 20:
                        text_chunks.append(chunk)
    except UnicodeDecodeError:
        # Try with different encoding if utf-8 fails
        try:
            with open(file_path, 'r', encoding='cp949') as file:  # Common Korean encoding
                content = file.read()
                
                if content:
                    chunks = chunk_text(content)
                    for chunk in chunks:
                        if len(chunk.strip()) > 20:
                            text_chunks.append(chunk)
        except Exception as e:
            print(f"Error extracting text from TXT with cp949 encoding: {e}")
    except Exception as e:
        print(f"Error extracting text from TXT: {e}")
    
    return text_chunks

def extract_text_from_csv(file_path: str) -> List[str]:
    """
    Extract text from CSV files with appropriate formatting for search.
    Treats CSV files similar to Excel but with a single sheet structure.
    """
    text_chunks = []
    
    try:
        # 다양한 인코딩 시도
        encodings = ['utf-8', 'cp949', 'euc-kr']
        df = None
        
        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                print(f"CSV 파일 '{os.path.basename(file_path)}' {encoding} 인코딩으로 성공적으로 읽음")
                break
            except UnicodeDecodeError:
                continue
        
        # 모든 인코딩으로 실패한 경우 기본 시스템 인코딩 시도
        if df is None:
            df = pd.read_csv(file_path)
            print(f"CSV 파일 '{os.path.basename(file_path)}' 시스템 기본 인코딩으로 읽음")
        
        # Get the filename for reference
        filename = os.path.basename(file_path)
        
        # Check if this might be a procedure guide CSV
        is_procedure_guide = any(keyword in filename for keyword in ['업무 안내', '업무_안내', '업무안내', '가이드', '매뉴얼', '절차'])
        
        # For procedure guides, process each row as separate knowledge chunk
        if is_procedure_guide:
            for idx, row in df.iterrows():
                # Skip empty rows
                if row.isna().all():
                    continue
                
                # Format as a structured knowledge snippet
                row_chunk_text = f"[업무 안내] "
                
                # Add all non-NA fields with their column names
                for col_name, value in row.items():
                    if not pd.isna(value) and str(value).strip():
                        # Clean the value
                        clean_value = str(value).strip().replace('\n', ' ')
                        row_chunk_text += f"{col_name}: {clean_value} | "
                
                # Remove trailing separator and add source reference
                row_chunk_text = row_chunk_text.rstrip(" | ")
                row_chunk_text += f"\n출처: {filename}"
                
                if len(row_chunk_text.strip()) > 20:  # Only keep meaningful chunks
                    text_chunks.append(row_chunk_text)
        else:
            # For regular CSV files, convert to plain text format
            csv_text = f"--- CSV: {filename} ---\n"
            
            # Add header row
            header_row = " | ".join([str(col) for col in df.columns])
            csv_text += header_row + "\n"
            
            # Add separator
            csv_text += "-" * len(header_row) + "\n"
            
            # Add data rows
            for _, row in df.iterrows():
                row_text = " | ".join([str(val) if not pd.isna(val) else "" for val in row])
                csv_text += row_text + "\n"
            
            # Split the CSV text into chunks
            from_function_chunks = chunk_text(csv_text)
            for chunk in from_function_chunks:
                if len(chunk.strip()) > 20:
                    text_chunks.append(chunk)
    
    except Exception as e:
        print(f"Error extracting text from CSV: {e}")
    
    return text_chunks

def extract_text_from_excel(file_path: str) -> List[str]:
    """
    Extract text from Excel files with special handling for procedure guide sheets.
    Specifically processes sheets like '절차_안내' (procedure guide) to create
    better chunks for RAG retrieval.
    """
    if pd is None:
        raise ImportError("pandas and openpyxl are required for Excel processing")
    
    text_chunks = []
    filename = os.path.basename(file_path)
    
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        procedure_sheet_found = False
        
        # First pass - look for procedure guide sheets
        for sheet_name in excel_file.sheet_names:
            # Check if this is a procedure guide sheet
            if "절차_안내" in sheet_name or "절차안내" in sheet_name:
                procedure_sheet_found = True
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Process each row in the procedure guide as a separate chunk for better retrieval
                for idx, row in df.iterrows():
                    # Skip empty rows
                    if row.isna().all():
                        continue
                    
                    # Format as a structured knowledge snippet
                    row_chunk_text = f"[업무 안내] "
                    
                    # Add all non-NA fields with their column names
                    for col_name, value in row.items():
                        if not pd.isna(value) and str(value).strip():
                            # Clean the value
                            clean_value = str(value).strip().replace('\n', ' ')
                            row_chunk_text += f"{col_name}: {clean_value} | "
                    
                    # Remove trailing separator and add source reference
                    row_chunk_text = row_chunk_text.rstrip(" | ")
                    row_chunk_text += f"\n출처: {filename} - {sheet_name} 시트"
                    
                    if len(row_chunk_text.strip()) > 20:  # Only keep meaningful chunks
                        text_chunks.append(row_chunk_text)
        
        # Second pass - process all sheets normally if no procedure guide was found,
        # or process non-procedure sheets in addition to procedure guides
        if not procedure_sheet_found or len(excel_file.sheet_names) > 1:
            for sheet_name in excel_file.sheet_names:
                # Skip re-processing procedure guide sheets that were already processed
                if procedure_sheet_found and ("절차_안내" in sheet_name or "절차안내" in sheet_name):
                    continue
                
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to string representations
                sheet_text = f"--- Sheet: {sheet_name} ---\n"
                
                # Add header row
                header_row = " | ".join([str(col) for col in df.columns])
                sheet_text += header_row + "\n"
                
                # Add separator
                sheet_text += "-" * len(header_row) + "\n"
                
                # Add data rows
                for _, row in df.iterrows():
                    row_text = " | ".join([str(val) if not pd.isna(val) else "" for val in row])
                    sheet_text += row_text + "\n"
                
                # Add empty line between sheets
                sheet_text += "\n"
                
                # Split the sheet text into chunks manually using the function
                from_function_chunks = chunk_text(sheet_text)
                for chunk in from_function_chunks:
                    if len(chunk.strip()) > 20:
                        text_chunks.append(chunk)
                        
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
    
    return text_chunks

def process_text(text: str, doc_id: str, filename: str, file_type: str = "txt", additional_metadata: dict = None) -> List[Dict[str, Any]]:
    """
    주어진 텍스트를 처리하고 청크로 분할하여 벡터 DB 형식으로 반환합니다.
    
    Args:
        text: 처리할 텍스트 내용
        doc_id: 문서 고유 ID
        filename: 원본 파일명
        file_type: 파일 형식 (기본값: txt)
        additional_metadata: 추가할 메타데이터 (선택 사항)
        
    Returns:
        메타데이터가 포함된 청크 목록
    """
    chunks = []
    
    # 기본 메타데이터 설정
    if additional_metadata is None:
        additional_metadata = {}
    
    # 텍스트를 청크로 분할
    raw_chunks = chunk_text(text)
    
    # 각 청크에 메타데이터 추가
    for i, chunk in enumerate(raw_chunks):
        # 기본 메타데이터
        metadata = {
            "filename": filename,
            "file_type": file_type,
            "chunk_index": i,
            "source": filename
        }
        
        # 추가 메타데이터 병합
        metadata.update(additional_metadata)
        
        chunks.append({
            "doc_id": doc_id,
            "chunk_id": f"{doc_id}-{i}",
            "text": chunk,
            "metadata": metadata
        })
    
    return chunks

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
    """
    Split text into overlapping chunks of approximately chunk_size characters
    
    Args:
        text: The text to split into chunks
        chunk_size: Target size of each chunk (500 tokens/characters by default)
        overlap: Number of characters to overlap between chunks
        
    Returns:
        List of text chunks
    """
    # Clean the text
    text = re.sub(r'\s+', ' ', text).strip()
    
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        # Find the end of the chunk
        end = start + chunk_size
        
        if end >= len(text):
            chunks.append(text[start:])
            break
        
        # Try to end at a sentence or paragraph boundary
        sentence_end = text.rfind('. ', start, end)
        para_end = text.rfind('\n', start, end)
        
        # Choose the latest boundary found within the range
        if sentence_end > start + chunk_size // 2:
            end = sentence_end + 2  # Include the period and space
        elif para_end > start + chunk_size // 2:
            end = para_end + 1  # Include the newline
        
        chunks.append(text[start:end])
        start = end - overlap  # Create overlap with the next chunk
    
    return chunks
